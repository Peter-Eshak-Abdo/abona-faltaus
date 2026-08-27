# 1. سكربت استخراج البيانات من ملفات البوربوينت وتحويلها إلى JSON[cite: 1]
# pip install python-pptx

import os
import json
from pptx import Presentation

ROOT_DIR = r"P:\قداس ملايكة\قداس ملايكةةةةةة\St.Mary Elnozha Liturgy Powerpoint"
OUTPUT_JSON = "./coptic_liturgy_dataset.json"

def extract_text_from_pptx(filepath):
    text_content = []
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if (hasattr(shape, "text") and shape.text.strip()):
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_content.append("\n".join(slide_text))
    except Exception:
        pass
    return text_content

def process_directory(root_path):
    dataset = []
    for dirpath, _, filenames in os.walk(root_path):
        for filename in filenames:
            if filename.lower().endswith('.pptx'):
                filepath = os.path.join(dirpath, filename)
                rel_path = os.path.relpath(filepath, root_path)
                path_parts = rel_path.split(os.sep)

                category = path_parts[0] if len(path_parts) > 1 else "liturgy"
                season = path_parts[1] if len(path_parts) > 2 else "annual"

                slides_text = extract_text_from_pptx(filepath)

                for idx, text in enumerate(slides_text):
                    item = {
                        "instruction": f"أكمل النص الطقسي التالي من {filename.replace('.pptx', '')}",
                        "input": "",
                        "output": text,
                        "category": category,
                        "season": season
                    }
                    dataset.append(item)
    return dataset

if __name__ == "__main__":
    data = process_directory(ROOT_DIR)
    with open(OUTPUT_JSON, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=4)
