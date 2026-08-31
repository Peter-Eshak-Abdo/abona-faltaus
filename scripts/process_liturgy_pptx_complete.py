import os
import json
import re
import urllib.parse
from collections import defaultdict
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = r"St.Mary Elnozha Liturgy Powerpoint"
OUTPUT_LITURGIES = r"lib/liturgies/data/full_liturgies_data.json"
OUTPUT_TASBEHA = r"lib/tasbeha/data/full_tasbeha_data.json"
OUTPUT_GRAPH = r"data/liturgy_hyperlink_graph.json"

def clean_txt(s):
    if not s:
        return ""
    # normalize spaces
    s = re.sub(r'[ \t]+', ' ', s)
    return s.strip()

def is_coptic_unicode(text):
    return any('\u2C80' <= ch <= '\u2CFF' or '\u03E2' <= ch <= '\u03EF' for ch in text)

def parse_slide_content(slide):
    title = ""
    arabic_texts = []
    coptic_texts = []
    coptic_arab_texts = []
    english_texts = []
    hyperlinks = []
    
    # Extract from shapes and tables
    for shape in slide.shapes:
        # Check Shape Hyperlinks (Buttons)
        try:
            if shape.click_action and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                hyperlinks.append({
                    "text": clean_txt(shape.text) if hasattr(shape, "text") and shape.text else "زر انتقال",
                    "target": shape.click_action.hyperlink.address
                })
        except Exception:
            pass

        # Text Frame
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p_text = clean_txt(p.text)
                if not p_text:
                    continue
                # Check run level hyperlinks
                for r in p.runs:
                    if r.hyperlink and r.hyperlink.address:
                        hyperlinks.append({
                            "text": clean_txt(r.text),
                            "target": r.hyperlink.address
                        })
                
                # Check if it's title
                if (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and "title" in shape.name.lower()) or (not title and len(p_text) < 60 and ("لحن" in p_text or "أوشية" in p_text or "ذكصولوجية" in p_text or "مرد" in p_text or "البولس" in p_text or "الإبركسيس" in p_text or "الكاثوليكون" in p_text or "الإنجيل" in p_text or "المزمور" in p_text or "الهوس" in p_text or "المجمع" in p_text or "صلاة" in p_text)):
                    if not title:
                        title = p_text
                        continue

                # Classify line
                if is_coptic_unicode(p_text):
                    coptic_texts.append(p_text)
                else:
                    arabic_texts.append(p_text)

        # Table (Col 1: Coptic, Col 2: Arabic, Col 3: Coptic-Arab)
        if shape.has_table:
            for row in shape.table.rows:
                row_cells = [clean_txt(c.text) for c in row.cells]
                non_empty = [c for c in row_cells if c]
                if len(non_empty) == 3:
                    coptic_texts.append(row_cells[0])
                    arabic_texts.append(row_cells[1])
                    coptic_arab_texts.append(row_cells[2])
                elif len(non_empty) == 2:
                    if is_coptic_unicode(non_empty[0]) or any(c.isascii() for c in non_empty[0] if c.isalpha()):
                        coptic_texts.append(non_empty[0])
                        arabic_texts.append(non_empty[1])
                    else:
                        arabic_texts.append(non_empty[0])
                        coptic_arab_texts.append(non_empty[1])
                elif len(non_empty) == 1:
                    arabic_texts.append(non_empty[0])

    return {
        "title": title,
        "arabic": "\n".join(arabic_texts).strip(),
        "coptic": "\n".join(coptic_texts).strip() if coptic_texts else None,
        "coptic_arabic": "\n".join(coptic_arab_texts).strip() if coptic_arab_texts else None,
        "english": "\n".join(english_texts).strip() if english_texts else None,
        "hyperlinks": hyperlinks
    }

def process_presentation(file_path):
    rel_path = os.path.relpath(file_path, BASE_DIR).replace("\\", "/")
    try:
        prs = pptx.Presentation(file_path)
    except Exception as e:
        return None

    raw_slides = []
    for idx, slide in enumerate(prs.slides):
        parsed = parse_slide_content(slide)
        parsed["slide_number"] = idx + 1
        raw_slides.append(parsed)

    # Intelligent Aggregation of multi-slide hymns / continuous readings
    sections = []
    current_hymn = None

    for slide in raw_slides:
        slide_title = slide["title"]
        has_content = bool(slide["arabic"] or slide["coptic"] or slide["coptic_arabic"])
        
        # If slide has a new title, start a new Section/Hymn
        if slide_title:
            if current_hymn and current_hymn["verses"]:
                sections.append(current_hymn)
            
            # Determine speaker & type
            speaker = "all"
            t_lower = slide_title.lower()
            if "كاهن" in t_lower or "أوشية" in t_lower or "صلاة" in t_lower or "إفلوجيمينوس" in t_lower:
                speaker = "priest"
            elif "شماس" in t_lower or "بروسيفكساستي" in t_lower or "طوبه هينا" in t_lower or "أنصتوا" in t_lower:
                speaker = "deacon"
            elif "شعب" in t_lower or "مرد" in t_lower or "لحن" in t_lower or "أرباع" in t_lower or "ذكصولوجية" in t_lower or "هوس" in t_lower:
                speaker = "people"
            elif "بولس" in t_lower or "كاثوليكون" in t_lower or "إبركسيس" in t_lower or "إنجيل" in t_lower or "سنكسار" in t_lower:
                speaker = "reader"

            stype = "hymn"
            if "أوشية" in t_lower or "صلاة" in t_lower:
                stype = "litany"
            elif "قسمة" in t_lower or "قسم" in t_lower:
                stype = "fraction"
            elif "قانون الإيمان" in t_lower:
                stype = "creed"
            elif "توزيع" in t_lower:
                stype = "communion"

            current_hymn = {
                "id": f"{rel_path}_s{slide['slide_number']}",
                "title": {
                    "arabic": slide_title,
                    "coptic": "",
                    "english": ""
                },
                "speaker": speaker,
                "type": stype,
                "verses": [],
                "hyperlinks": slide["hyperlinks"],
                "source_file": rel_path,
                "start_slide": slide["slide_number"]
            }

            if has_content:
                current_hymn["verses"].append({
                    "arabic": slide["arabic"],
                    "coptic": slide["coptic"],
                    "coptic_arabic": slide["coptic_arabic"],
                    "english": slide["english"],
                    "slide_number": slide["slide_number"]
                })
        else:
            # Continuing current hymn/reading (multi-slide hymn verses)
            if current_hymn is not None:
                if has_content:
                    current_hymn["verses"].append({
                        "arabic": slide["arabic"],
                        "coptic": slide["coptic"],
                        "coptic_arabic": slide["coptic_arabic"],
                        "english": slide["english"],
                        "slide_number": slide["slide_number"]
                    })
                if slide["hyperlinks"]:
                    current_hymn["hyperlinks"].extend(slide["hyperlinks"])
            elif has_content:
                current_hymn = {
                    "id": f"{rel_path}_s{slide['slide_number']}",
                    "title": {
                        "arabic": os.path.splitext(os.path.basename(file_path))[0],
                        "coptic": "",
                        "english": ""
                    },
                    "speaker": "all",
                    "type": "hymn",
                    "verses": [{
                        "arabic": slide["arabic"],
                        "coptic": slide["coptic"],
                        "coptic_arabic": slide["coptic_arabic"],
                        "english": slide["english"],
                        "slide_number": slide["slide_number"]
                    }],
                    "hyperlinks": slide["hyperlinks"],
                    "source_file": rel_path,
                    "start_slide": slide["slide_number"]
                }

    if current_hymn and current_hymn["verses"]:
        sections.append(current_hymn)

    return {
        "file": rel_path,
        "title": os.path.splitext(os.path.basename(file_path))[0],
        "total_slides": len(prs.slides),
        "sections": sections
    }

def main():
    print("Starting comprehensive extraction & grouping...")
    all_presentations = []
    hyperlink_graph = defaultdict(list)

    for root, dirs, files in os.walk(BASE_DIR):
        for f in files:
            if f.lower().endswith('.pptx') and not f.startswith('~$'):
                full_p = os.path.join(root, f)
                res = process_presentation(full_p)
                if res and res["sections"]:
                    all_presentations.append(res)
                    for sec in res["sections"]:
                        for hl in sec.get("hyperlinks", []):
                            target = hl["target"]
                            hyperlink_graph[res["file"]].append({
                                "from_section": sec["id"],
                                "label": hl["text"],
                                "target": target
                            })

    print(f"Successfully processed {len(all_presentations)} presentations!")
    
    # Save hyperlink graph
    os.makedirs(os.path.dirname(OUTPUT_GRAPH), exist_ok=True)
    with open(OUTPUT_GRAPH, "w", encoding="utf-8") as f:
        json.dump(hyperlink_graph, f, ensure_ascii=False, indent=2)
    print(f"Hyperlink graph saved to {OUTPUT_GRAPH}")

    # Build Liturgies and Tasbeha structured JSON collections
    liturgy_groups = defaultdict(list)
    tasbeha_groups = defaultdict(list)

    for p in all_presentations:
        fpath = p["file"]
        sec_count = len(p["sections"])
        v_count = sum(len(s["verses"]) for s in p["sections"])
        
        # Categorize by Folder
        top_folder = fpath.split('/')[0] if '/' in fpath else fpath
        
        group_item = {
            "id": p["file"].replace("/", "_").replace(".pptx", ""),
            "title": {
                "arabic": p["title"],
                "coptic": "",
                "english": ""
            },
            "badge": f"{p['total_slides']} شريحة ({v_count} ربع/فقرة)",
            "source_file": p["file"],
            "sections": p["sections"]
        }

        if "تسبحة" in fpath or "25 تسبحة" in fpath or "26 تسبحة" in fpath:
            tasbeha_groups[top_folder].append(group_item)
        else:
            liturgy_groups[top_folder].append(group_item)

    # Format Liturgy Documents
    liturgies_docs = []
    for top_cat, grps in liturgy_groups.items():
        liturgies_docs.append({
            "id": top_cat.replace(" ", "-"),
            "slug": top_cat.replace(" ", "-"),
            "title": {
                "arabic": top_cat,
                "coptic": "",
                "english": top_cat
            },
            "subtitle": f"مجموعة صلوات وطقوس {top_cat}",
            "description": f"تحتوي على {len(grps)} ملف صلاة وقراءات من عروض الباوربوينت المنظمة بدقة.",
            "iconName": "FaChurch",
            "accentColor": "amber",
            "groups": grps
        })

    # Format Tasbeha Documents
    tasbeha_docs = []
    for top_cat, grps in tasbeha_groups.items():
        tasbeha_docs.append({
            "id": top_cat.replace(" ", "-"),
            "slug": top_cat.replace(" ", "-"),
            "title": {
                "arabic": top_cat,
                "coptic": "",
                "english": top_cat
            },
            "subtitle": f"تسابيح ومدائح وإبصالمودية {top_cat}",
            "description": f"تحتوي على {len(grps)} ملف تسبحة وهوسات ومدائح منظمة بالأرباع.",
            "iconName": "FaMusic",
            "accentColor": "emerald",
            "groups": grps
        })

    # Save to lib/liturgies and lib/tasbeha
    os.makedirs(os.path.dirname(OUTPUT_LITURGIES), exist_ok=True)
    with open(OUTPUT_LITURGIES, "w", encoding="utf-8") as f:
        json.dump(liturgies_docs, f, ensure_ascii=False, indent=2)
    print(f"Saved {len(liturgies_docs)} Liturgy documents to {OUTPUT_LITURGIES}")

    os.makedirs(os.path.dirname(OUTPUT_TASBEHA), exist_ok=True)
    with open(OUTPUT_TASBEHA, "w", encoding="utf-8") as f:
        json.dump(tasbeha_docs, f, ensure_ascii=False, indent=2)
    print(f"Saved {len(tasbeha_docs)} Tasbeha documents to {OUTPUT_TASBEHA}")

if __name__ == "__main__":
    main()
